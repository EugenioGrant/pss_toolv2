import streamlit as st
import pandas as pd
import numpy as np
from pandas import DataFrame
import base64
import os
from PIL import Image
import requests
from io import BytesIO
from pptx import Presentation
from pptx.chart.data import XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.util import Pt
from pptx.enum.chart import XL_MARKER_STYLE
import datetime

#-------------------------------------------Functions-------------------------------------------#
def file_selector(folder_path='.'): # Function to select all files in a specific local folder
    filenames = os.listdir(folder_path)
    selected_filename = st.selectbox('Select a file', filenames)
    return os.path.join(folder_path, selected_filename)

def csv_downloader(data): # Function to download only CSV files
    csv_file = data.to_csv(index=False)
    b64 = base64.b64encode(csv_file.encode()).decode()
    new_filename = "Output_data.csv"
    st.markdown("#### Download File ###")
    href = f'<a href="data:file/csv;base64,{b64}" download="{new_filename}">Click Here!!</a>'
    st.markdown(href, unsafe_allow_html=True)

class FileDownloader(object): # Class to enable general file download online

    def __init__(self, data,filename='PSSReport',file_ext='pptx'):
        super(FileDownloader, self).__init__()
        self.data = data
        self.filename = filename
        self.file_ext = file_ext

    def download(self):
        b64 = base64.b64encode(self.data).decode()
        new_filename = "{}.{}".format(self.filename,self.file_ext)
        st.markdown("#### Download File ###")
        href = f'<a href="data:file/{self.file_ext};base64,{b64}" download="{new_filename}">Click Here!!</a>'
        st.markdown(href,unsafe_allow_html=True)

def applyFunc(s): #check if time of respondent is lower then cutoff time
    if s < MINIMUM_TIME * median_time:
        return 0
    else:
        return 1

def price_sensitivity_meter(df, number, interpolate=False): # PSS processing function
    # convert data from wide to long
    # calculate frequency of each price for each group
    df1 = (df[['Too Cheap', 'Cheap', 'Expensive', 'Too Expensive']]
           .unstack()
           .reset_index()
           .rename(columns={'level_0': 'label', 0: 'prices'})[['label', 'prices']]
           .groupby(['label', 'prices'])
           .size()
           .reset_index()
           .rename(columns={0: 'frequency'})
           )

    # calculate cumsum percentages
    df1['cumsum'] = df1.groupby(['label'])['frequency'].cumsum()
    df1['sum'] = df1.groupby(['label'])['frequency'].transform('sum')
    df1['percentage'] = 100 * df1['cumsum'] / df1['sum']
    # convert data from long back to wide
    df2 = df1.pivot_table('percentage', 'prices', 'label')

    # take linear values in missing values
    if interpolate:
        df3 = df2.interpolate().fillna(0)
        df3['Too Cheap'] = 100 - df3['Too Cheap']
        df3['Cheap'] = 100 - df3['Cheap']

    # forward fill
    else:
        df3 = df2.ffill().fillna(0)

        df3['Too Cheap'] = 100 - df3['Too Cheap']
        df3['Cheap'] = 100 - df3['Cheap']

    df3['optimal_diff'] = (df3['Too Cheap'] - df3['Too Expensive'])
    df3['left_diff'] = (df3['Too Cheap'] - df3['Expensive'])
    df3['right_diff'] = (df3['Too Expensive'] - df3['Cheap'])
    optimal = df3[df3['optimal_diff'] <= 0].index[0]
    lower_bound = df3[df3['left_diff'] <= 0].index[0]
    upper_bound = df3[df3['right_diff'] >= 0].index[0]

    df3['Product'] = number

    return df3, optimal, lower_bound, upper_bound

#-------------------------------------------Initiate sidebar menu--------------------------------------------------#
menu = ["Start","Data Cleaning", "Data Processing", "PPT Report Creation"]

st.sidebar.title("Navigation Menu")
choice = st.sidebar.selectbox("Options", menu)

#-------------------------------------------Code for each tab in web app-------------------------------------------#
if choice == "Start":
    response = requests.get("https://skimgroup.com/app/uploads/2017/02/facebook.png")
    img = Image.open(BytesIO(response.content))
    new_img=img.resize((800,400))
    st.image(new_img, output_format="PNG")
    st.title("SKIM PSS Tool")

    #skimlogo
    st.subheader("**Welcome to the SKIM PSS Tool**")
    st.subheader("Instructions: ")
    st.subheader("Please navigate using the Menu sidebar to the left. ")
    st.subheader("This tool was designed to simplify the data cleaning, data processing and reporting phases of the VW price check studies.")
    st.subheader("Please write to: **m.villegas@skimgroup.com** for any questions, issues or suggestions.")

    st.title("Manual")

    st.subheader("Please click the link below for the Tool manual:")
    href = f'<a href="https://skimgroup.sharepoint.com/:p:/r/sites/Teams/LATAM/Tools%20LATAM/PSS_Tool/PSS_Tool_Manual.pptx?d=wd00e22843f9b4855b272eae6bfc382d5&csf=1&web=1&e=h2pGYj" >Go to Manual</a>'
    st.markdown(href, unsafe_allow_html=True)
    href = f'<a href="https://skimgroup.sharepoint.com/:f:/r/sites/Teams/LATAM/Tools%20LATAM/PSS_Tool?csf=1&web=1&e=tV5wyX" >Go to Tool Folder</a>'
    st.markdown(href, unsafe_allow_html=True)

elif choice == "Data Cleaning":
    st.title("Data Cleaning Tool")
    st.subheader("1. Data Cleaning Settings (General)")
    PRODUCT_NUMBER = st.number_input("Number of Products", min_value=1, max_value=100, value=12, step=1)
    MINIMUM_TIME = st.number_input("Minimum Time (X*median)", min_value=0.10, max_value=0.90, value=0.3,step=1/10.,format="%.2f")

    st.subheader("2. Data Cleaning Settings (Clean Price Awareness Outliers)")
    MIN_PERCENTILE = st.number_input("Minimum Percentile", min_value=0.00, max_value=0.25, value=0.005,step=1/1000.,format="%.3f")
    MAX_PERCENTILE = st.number_input("Maximum Percentile", min_value=0.75, max_value=1.00, value=0.995, step=1 / 1000.,format="%.3f")

    st.subheader("3. Upload & Process Raw Data file")
    data_file = st.file_uploader("Please upload CSV", type=['csv'])
    if data_file is not None:
        st.success("File uploaded correctly")

    if st.button("Process"):
        if data_file is not None:
            file_details = {"Filename": data_file.name, "FileType": data_file.type, "FileSize": data_file.size}
            df = pd.read_csv(data_file)

            # -------------------------- Initiate Blank arrays --------------------------
            Awareness2 = []  # initiating blank arrays
            cheap2 = []
            too_cheap2 = []
            expensive2 = []
            too_expensive2 = []
            extra_cheap2 = []
            extra_expensive2 = []
            used = []

            # -------------------------- Append metrics for each product --------------------------
            for i in range(PRODUCT_NUMBER):
                Awareness = 'Q1' + '.' + str(i + 1)
                Awareness2.append(Awareness)
                cheap = 'Q2a' + '.' + str(i + 1)
                cheap2.append(cheap)
                too_cheap = 'Q2b' + '.' + str(i + 1)
                too_cheap2.append(too_cheap)
                expensive = 'Q2c' + '.' + str(i + 1)
                expensive2.append(expensive)
                too_expensive = 'Q2d' + '.' + str(i + 1)
                too_expensive2.append(too_expensive)
                extra_cheap = 'Q3' + '.' + str(i + 1)
                extra_cheap2.append(extra_cheap)
                extra_expensive = 'Q4' + '.' + str(i + 1)
                extra_expensive2.append(extra_expensive)

            # -------------------------- New dataframe with only necessary columns --------------------------
            df2 = df[['sys_RespNum', 'sys_StartTimeStamp', 'sys_EndTimeStamp']]
            study_time2 = df.loc[:, 'sys_EndTimeStamp'] - df.loc[:, 'sys_StartTimeStamp']
            df2['study_time'] = study_time2
            median_time = df2['study_time'].median()
            df2['used'] = df2['study_time'].apply(applyFunc)

            df2[Awareness2] = df[Awareness2]    #Q1
            df2[cheap2] = df[cheap2]            #Q2a
            df2[too_cheap2] = df[too_cheap2]    #Q2b
            df2[expensive2] = df[expensive2]    #Q2c
            df2[too_expensive2] = df[too_expensive2]    #Q2d
            df2[extra_cheap2] = df[extra_cheap2]        #Q3
            df2[extra_expensive2] = df[extra_expensive2] #Q4
            df2 = df2.fillna(0)

            # -------------------------- Initiate loop variables for final data dataframe --------------------------
            new_resp = []
            respondent = []
            Product = []
            used = []
            Price_Awareness = []
            Q2b = []
            Q2a = []
            Q2c = []
            Q2d = []
            Q3 = []
            Q4 = []
            cleaned = 0
            # st.dataframe(df2)
            # -------------------------- Loop to organize data --------------------------
            for i in range(PRODUCT_NUMBER):
                for k in range(len(df2)):
                    if df2.iloc[k, 5 + PRODUCT_NUMBER + i] != 0:
                        new_resp.append(df2.iloc[k, 0] * 1000 + (i + 1))
                        Product.append(i + 1)
                        Q2b.append(df2.iloc[k, 5 + 2 * PRODUCT_NUMBER + i])
                        Q2a.append(df2.iloc[k, 5 + PRODUCT_NUMBER + i])
                        Q2c.append(df2.iloc[k, 5 + 3 * PRODUCT_NUMBER + i])
                        Q2d.append(df2.iloc[k, 5 + 4 * PRODUCT_NUMBER + i])
                        Q3.append(df2.iloc[k, 5 + 5 * PRODUCT_NUMBER + i])
                        Q4.append(df2.iloc[k, 5 + 6 * PRODUCT_NUMBER + i])
                        Price_Awareness.append(df2.iloc[k, 5 + i])
                        respondent.append(df2.iloc[k, 0])
                        used.append(df2.iloc[k, 4])

            # -------------------------- consolidate metrics into single array --------------------------
            consolidated_list = [new_resp, Product, Q2b, Q2a, Q2c, Q2d, Q3, Q4, Price_Awareness, respondent, used]

            # -------------------------- convert array into dataframe --------------------------
            export = DataFrame(consolidated_list).transpose()
            export.columns = ['new resp', 'Product', 'Q2b', 'Q2a', 'Q2c', 'Q2d', 'Q3', 'Q4', 'Price Awareness',
                              'Respondent', 'USED']
            # export5 = export.sort_values('Price Awareness')
            # st.dataframe(export5)
            used2 = []
            new_df2 = []
            top_percentile_list = []
            min_percentile_list = []
            # -------------------------- Clean data based on price awareness --------------------------

            for x in range(PRODUCT_NUMBER):
                current = x + 1
                new_df = export.loc[export['Product'] == current]
                new_df = new_df.reset_index()
                top_percentile_list.append(np.percentile(new_df['Price Awareness'], MAX_PERCENTILE * 100))
                min_percentile_list.append(np.percentile(new_df['Price Awareness'], MIN_PERCENTILE * 100))

            export = export.sort_values('Price Awareness')
            export = export.reset_index()

            for respondent in df['sys_RespNum']:
                new_df = export.loc[export['Respondent'] == respondent]
                new_df = new_df.reset_index()
                for k in range(len(new_df)):
                    if new_df.loc[k, 'USED'] == 0:
                        cleaned = cleaned + 1
                        for j in range(len(new_df)):
                            used2.append(0)
                        break
                    elif new_df.loc[k, 'Price Awareness'] > top_percentile_list[int(new_df.loc[k, 'Product'] - 1)] or \
                            new_df.loc[k, 'Price Awareness'] < min_percentile_list[int(new_df.loc[k, 'Product'] - 1)]:
                        cleaned = cleaned + 1
                        for w in range(len(new_df)):
                            used2.append(0)
                        break
                    else:
                        # if new_df.loc[k, 'USED'] > 0:
                        for k in range(len(new_df)):
                            used2.append(1)
                        break

            export = export.sort_values('Respondent')
            export['USED FINAL'] = np.array(used2)

            Q4_col = export['Q4']
            Q4_col.replace(to_replace = 0, value = 1, inplace=True)

            export2 = export.drop(['Price Awareness', 'Respondent', 'USED', 'index'], 1)
            export2 = export2.loc[export2['USED FINAL'] == 1]
            export2 = export2.drop(['USED FINAL'], 1)
            export = export.drop(['index'], 1)

            median_price = []
            min_price = []
            max_price = []
            current_product = []
            for w in range(PRODUCT_NUMBER):
                current = w + 1
                export3 = export.loc[export['USED FINAL'] == 1]
                new_df = export3.loc[export3['Product'] == current]
                new_df = new_df.reset_index()
                median_price.append(new_df['Price Awareness'].median())
                min_price.append(min(new_df['Price Awareness']))
                max_price.append(max(new_df['Price Awareness']))
                current_product.append(current)

            consolidated_list2 = [current_product, median_price, min_price, max_price]

            # -------------------------- convert array into dataframe --------------------------
            export_price = DataFrame(consolidated_list2).transpose()
            export_price.columns = ['Product','Median', 'Minimum', 'Maximum']
            
            export_price.index = [""]*len(export_price)

            st.subheader('4. Check Median, minimum and maximum Price Awareness')
            st.table(export_price)
    
            st.subheader('5. Cleaned respondents: {}, of Total sample: {} - Remaining Sample: {}'.format(cleaned, len(df), len(df) - cleaned))
            st.success("Data processed successfully")

            st.subheader("Download processed Data File (VW data)")
            csv_downloader(export2)
            st.subheader("Download processed Data File (Long VW data *Includes USED column)")
            csv_downloader(export)

elif choice == "Data Processing":
    st.title("Data Processing Tool")
    st.write("if the Excel simulator is needed, please go to the folder below, download the excel tool, input the data from the previous step and run the macro to process the data.")
    st.write("*For most clients we send the excel-based simulator to them.")
    href = f'<a href="https://skimgroup.sharepoint.com/:f:/r/sites/Teams/LATAM/Tools%20LATAM/PSS_Tool?csf=1&web=1&e=tV5wyX" >Go to Tool Folder</a>'
    st.markdown(href, unsafe_allow_html=True)

    st.subheader("1. Data Processing Settings (General)")
    col1, col2, col3 = st.beta_columns(3)

    with col1:
        COUNTRY = st.text_input('Country', 'Country')
    with col2:
        time_now = datetime.datetime.now()
        DATE = st.date_input("Date of FW", time_now)
        DATE2 = DATE.strftime("%x")
        # st.write(DATE2)
    with col3:
        PRODUCT_NUMBER = st.number_input("Number of Products", min_value=1, max_value=100, value=12, step=1)

    st.subheader("2. Upload & Process Cleaned Data file")
    data_file = st.file_uploader("Please upload CSV", type=['csv'])
    if data_file is not None:
        st.success("File uploaded correctly")

    names = []
    current_prices = []

    st.subheader("3. Include Product Information (Name & Cost)")
    SKUcol1, SKUcol2 = st.beta_columns(2)
    for i in range(PRODUCT_NUMBER):
        with SKUcol1:
            temp_name = st.text_input(f'Product: {i+1}', f'SKU{i+1}')
            names.append(temp_name)
        with SKUcol2:
            temp_cost = st.text_input(f'Product: {i + 1}', 15)
            current_prices.append(temp_cost)

    if st.button("Process"):
        if data_file is not None:
            database = pd.read_csv(data_file)
            df_final = pd.DataFrame(
                            columns=["ProductNumber", "values", "not cheap", "not expensive", "too cheap", "too expensive", "optimal",
                            "lower_bound", "upper_bound"])
            number_of_products = max(database.loc[:, 'Product'])

            latest_iteration = st.empty()
            bar = st.progress(0)

            for i in range(int(number_of_products)):
                current = i + 1
                latest_iteration.text(f'Processing Curve: {current}/{int(number_of_products)} - {int((100 / number_of_products) * current)}%')
                bar.progress(int((100 / number_of_products) * current))
                new_df = database.loc[database['Product'] == current]
                new_df = new_df.reset_index()

                prices = {'Too Cheap': new_df['Q2b'],
                          'Cheap': new_df['Q2a'],
                          'Expensive': new_df['Q2c'],
                          'Too Expensive': new_df['Q2d'],
                          }

                df = pd.DataFrame(prices)

                df3, optimal, lower_bound, upper_bound = price_sensitivity_meter(df, current)

                df3['optimal'] = optimal
                df3['lower_bound'] = lower_bound
                df3['upper_bound'] = upper_bound
                df3['Prices'] = df3.index

                df4 = df3[
                    ["Product", "Prices", "Cheap", "Expensive", "Too Cheap", "Too Expensive", "optimal", "lower_bound",
                     "upper_bound"]]
                df4 = df4.reset_index()

                df4["Cheap"] = 1 - df4["Cheap"].div(100).round(4)
                df4["Expensive"] = 1 - df4["Expensive"].div(100).round(4)
                df4["Too Cheap"] = df4["Too Cheap"].div(100).round(4)
                df4["Too Expensive"] = df4["Too Expensive"].div(100).round(4)

                df4.rename(columns={'Cheap': 'not cheap', 'Expensive': 'not expensive', 'Too Cheap': 'too cheap', 'Too Expensive': 'too expensive', 'Prices': 'values', 'Product': 'ProductNumber'}, inplace=True)
                df4.drop('prices', axis=1, inplace=True)

                df_final = df_final.append(df4)

            df_final['Date'] = DATE2
            df_final['Country'] = COUNTRY
            df_final['Restaurant'] = 'NA'
            new_df3 = pd.DataFrame()

            for i in range(PRODUCT_NUMBER):
                current = i+1
                df_final['ProductName'] = ''
                new_df2 = df_final.loc[df_final['ProductNumber'] == current]
                new_df2['ProductName'] = names[i]
                new_df2['Real_Price'] = current_prices[i]

                # x_values = new_df2['values'].to_numpy()
                # y_values = new_df2['not expensive'].to_numpy()
                # y_f = interp1d(x_values, y_values, 'linear')
                # NE_Real_Price = y_f(current_prices[i])

                x_values = new_df2['values'].to_numpy().astype(float)
                y_values = new_df2['not expensive'].to_numpy().astype(float)
                NE_Real_Price = np.interp(current_prices[i], x_values, y_values)

                new_df2['NE_RealPrice'] = NE_Real_Price

                new_df3 = new_df3.append(new_df2)

            df_processed_final = new_df3[
                ["Date", "Country", "Restaurant", "ProductNumber", "ProductName", "values", "too cheap", "not cheap", "not expensive", "too expensive","Real_Price","NE_RealPrice", "optimal", "lower_bound",
                 "upper_bound"]]

            st.subheader("Download processed Data File (VW Results)")
            csv_downloader(df_processed_final)

elif choice == "PPT Report Creation":
    st.title("PPT Report Creation Tool")

    st.subheader("1. Upload Processed VW Data file")
    data_file3 = st.file_uploader("Please upload CSV", type=['csv'])
    if data_file3 is not None:
        st.success("Vw data file uploaded correctly")

    check = st.checkbox("Include Trial, Revenue & Profit curves?")

    if check:
        st.subheader("1.2 Upload Processed Trial, Revenue & Profit Data file")
        data_file9 = st.file_uploader("Please upload Trial CSV", type=['csv'])
        if data_file9 is not None:
            st.success("Trial, Revenue & Profit Data uploaded correctly")

    st.subheader("2. Upload Powerpoint report layout file (pptx)")
    data_file2 = st.file_uploader("Please upload PPT", type=['pptx'])

    st.write('Template Settings:')

    first, last, last2 = st.beta_columns(3)
    legend_settings_names = ['English', 'Spanish']
    legend_settings = first.radio('Language Settings:', legend_settings_names)
    color_settings = last.radio('Chart Color Scheme Settings:', ['Mcdonalds', 'CMI'])

    if data_file2 is not None:
        st.success("PPT Layout file uploaded correctly")

    st.subheader("3. Upload Product Images (.png files - multiple select)")
    data_file10 = st.file_uploader("Please upload png images", type=['png'], accept_multiple_files=True)

    if len(data_file10) >0:
        st.success("Images uploaded correctly")

    temp = 0
    for i in range(0, len(data_file10)):
        for j in range(i + 1, len(data_file10)):
            current = data_file10[i].name.split(".")[0]
            next = data_file10[j].name.split(".")[0]
            if (int(current) > int(next)):
                temp = data_file10[i]
                data_file10[i] = data_file10[j]
                data_file10[j] = temp

    st.subheader("4. Run Report creation tool")

    prs = Presentation(data_file2)
    slide_layout = prs.slide_layouts[1]
    slide_layout2 = prs.slide_layouts[2]
    slide_layout3 = prs.slide_layouts[3]

    if st.button("Click to create and export report"):
        if data_file2 is not None:
            df = pd.read_csv(data_file3)

            if check:
                if data_file9 is not None:
                    df_trial = pd.read_csv(data_file9)

            number_of_products = max(df.loc[:, 'ProductNumber'])
            country = df.loc[0, 'Country']
            df['Date'] = pd.to_datetime(df.Date)
            date = df.loc[0, 'Date']
            # ----------------------------------------------------------------------------------------------------------------#
            #                  Loop through all products and create a slide for each of them
            latest_iteration = st.empty()
            bar = st.progress(0)

            for i in range(number_of_products):
                current = i + 1
                # st.write("Processing Curve: ", current, "/", number_of_products)
                latest_iteration.text(f'Processing Curve: {current}/{number_of_products} - {int((100/number_of_products)*current)}%')
                bar.progress(int((100/number_of_products)*current))

                new_df2 = df.loc[df['ProductNumber'] == current]
                new_df2 = new_df2.reset_index()
                # st.table(new_df2)
                # st.write(new_df2.loc[0, 'Real_Price'])
                Real_price = new_df2.loc[0, 'Real_Price']

                if color_settings == "CMI":
                    NE_Real_price = 1- new_df2.loc[len(new_df2)-1, 'not cheap']
                else:
                    NE_Real_price = new_df2.loc[0, 'NE_RealPrice']

                new_df = new_df2.sort_values('values')
                name1 = 'Not Expensive'
                name2 = 'Too Expensive'
                name3 = 'Not Cheap'
                name4 = 'Too Cheap'

                if legend_settings == 'Spanish':
                    if data_file9 is not None:
                        new_df_trial = df_trial.loc[df_trial['ProductNumber'] == current]
                        name1 = 'Caro'
                        name2 = 'Demasiado Caro'
                        name3 = 'Buen Precio'
                        name4 = 'Demasiado Barato'

                slide = prs.slides.add_slide(slide_layout)
                slide.placeholders[13].text = date.strftime("%b %Y") + ' - ' + new_df.loc[0, 'ProductName']  # Orange text in sub Header
                if legend_settings == 'Spanish':
                    slide.placeholders[0].text = country.upper() + ' - Analisis Van Westendorp'
                else:
                    slide.placeholders[0].text = country.upper() + ' - Van Westendorp Analysis'
                picture = slide.placeholders[14].insert_picture(data_file10[i])  # Product Image
                slide.placeholders[16].text = country
                slide.placeholders[17].text = date.strftime("%b %Y")

                chart_data = XyChartData()
                x_values = new_df['values'].to_numpy()

                if color_settings == "CMI":
                    y_values = 1 - new_df['not expensive'].to_numpy()  #Caro
                    y_values3 = 1 - new_df['not cheap'].to_numpy()  #Buen Precio
                else:
                    y_values = new_df['not expensive'].to_numpy()
                    y_values3 = new_df['not cheap'].to_numpy()

                y_values2 = new_df['too expensive'].to_numpy()
                y_values4 = new_df['too cheap'].to_numpy()

                minimum = min(x_values)
                maximum = max(x_values)
                units = round((maximum - minimum) / 14 + 0.5, 0)

                # ----------------------------------------------------------------------------------------------------------------#
                cd = chart_data.add_series(name1, number_format='$#,##0')
                for x, y in list(zip(x_values, y_values)):
                    cd.add_data_point(x, y, number_format=None)

                cd = chart_data.add_series(name2, number_format='$#,##0')
                for x, y in list(zip(x_values, y_values2)):
                    cd.add_data_point(x, y, number_format=None)

                cd = chart_data.add_series(name3, number_format='$#,##0')
                for x, y in list(zip(x_values, y_values3)):
                    cd.add_data_point(x, y, number_format=None)

                cd = chart_data.add_series(name4, number_format='$#,##0')
                for x, y in list(zip(x_values, y_values4)):
                    cd.add_data_point(x, y, number_format=None)

                cd = chart_data.add_series('Actual Price', number_format=None)
                cd.add_data_point(Real_price, 0, number_format=None)
                cd.add_data_point(Real_price, NE_Real_price, number_format=None)

                cd = chart_data.add_series('Actual', number_format=None)
                cd.add_data_point(Real_price, NE_Real_price, number_format=None)
                # ----------------------------------------------------------------------------------------------------------------#
                #                   Adjust format of chart

                graphic_frame = slide.placeholders[15].insert_chart(XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
                                                                    chart_data)
                chart = graphic_frame.chart
                chart.chart_style = 2  #
                chart.has_title = False  #
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.font.size = Pt(12)
                chart.legend.font.color.rgb = RGBColor(89, 89, 89)
                if color_settings == "CMI":
                    chart.font.name = 'Rawline'
                else:
                    chart.font.name = 'Raleway'

                category_axis = chart.category_axis
                category_axis.has_major_gridlines = True  #
                category_axis.tick_labels.font.size = Pt(9)
                category_axis.tick_labels.font.color.rgb = RGBColor(89, 89, 89)
                category_axis.maximum_scale = maximum  # Maximum
                category_axis.minimum_scale = minimum  # minimum
                category_axis.major_gridlines.format.line.color.rgb = RGBColor(217, 217, 217)
                category_axis.format.line.color.rgb = RGBColor(217, 217, 217)
                category_axis.major_unit = units
                # category_axis.axis_title.text_frame.text = 'Precio'

                value_axis = chart.value_axis
                value_axis.has_major_gridlines = False
                value_axis.tick_labels.font.size = Pt(9)
                value_axis.maximum_scale = 1.0  # Maximum
                value_axis.minimum_scale = 0.0  # minimum
                tick_labels = value_axis.tick_labels  # Tick_Labels Control class for the vertical axis tag of Chart
                tick_labels.number_format = '0%'  # Tag display style
                tick_labels.font.color.rgb = RGBColor(89, 89, 89)
                value_axis.major_gridlines.format.line.color.rgb = RGBColor(217, 217, 217)
                value_axis.format.line.color.rgb = RGBColor(217, 217, 217)
                # value_axis.axis_title.text_frame.text = 'Consumidores %'

                plot = chart.plots[0]
                series = plot.series[0]
                if color_settings == 'CMI':
                    series.format.line.color.rgb = RGBColor(0, 176, 240)
                else:
                    series.format.line.color.rgb = RGBColor(146, 208, 80)
                series.format.line.width = Pt(2.25)

                plot = chart.plots[0]
                series = plot.series[1]
                if color_settings == 'CMI':
                    series.format.line.color.rgb = RGBColor(255, 0, 0)
                else:
                    series.format.line.color.rgb = RGBColor(255, 0, 0)
                series.format.line.width = Pt(2.25)

                plot = chart.plots[0]
                series = plot.series[2]
                if color_settings == 'CMI':
                    series.format.line.color.rgb = RGBColor(198, 15, 83)
                else:
                    series.format.line.color.rgb = RGBColor(255, 255, 0)
                series.format.line.width = Pt(2.25)

                plot = chart.plots[0]
                series = plot.series[3]
                if color_settings == 'CMI':
                    series.format.line.color.rgb = RGBColor(142, 180, 227)
                else:
                    series.format.line.color.rgb = RGBColor(0, 176, 240)
                series.format.line.width = Pt(2.25)

                plot = chart.plots[0]
                series = plot.series[4]
                if color_settings == 'CMI':
                    series.format.line.color.rgb = RGBColor(198, 15, 83)
                else:
                    series.format.line.color.rgb = RGBColor(146, 208, 80)
                series.format.line.dash_style = 4

                plot = chart.plots[0]
                series = plot.series[5]
                if color_settings == 'CMI':
                    series.format.line.color.rgb = RGBColor(198, 15, 83)
                else:
                    series.format.line.color.rgb = RGBColor(146, 208, 80)
                series.marker.size = 8
                series.marker.format.fill.solid()
                series.marker.format.line.fill.background()
                if color_settings == 'CMI':
                    series.marker.format.fill.fore_color.rgb = RGBColor(198, 15, 83)
                else:
                    series.marker.format.fill.fore_color.rgb = RGBColor(146, 208, 80)
                series.marker.style = XL_MARKER_STYLE.SQUARE

                data_label = series.points[0].data_label
                data_label.position = XL_DATA_LABEL_POSITION.LEFT
                if color_settings == 'CMI':
                    data_label.font.color.rgb = RGBColor(198, 15, 83)
                else:
                    data_label.font.color.rgb = RGBColor(146, 208, 80)
                data_label.font.size = Pt(12)
                data_label.font.bold = True
                data_label.Number_Format = '0%'

                # ----------------------------------------------------------------------------------------------------------------#
                #               Add Not Expensive Slide for each product
                if color_settings == "CMI":
                    pass
                else:
                    slide = prs.slides.add_slide(slide_layout2)
                    slide.placeholders[13].text = date.strftime("%b %Y") + ' - ' + new_df.loc[
                        0, 'ProductName']  # Orange text in sub Header
                    if legend_settings == 'Spanish':
                        slide.placeholders[0].text = country.upper() + ' - Analisis Not Expensive'  # Blue Text in Header
                    else:
                        slide.placeholders[0].text = country.upper() + ' - Not Expensive Analysis'
                    picture = slide.placeholders[14].insert_picture(data_file10[i])  # Product Image
                    slide.placeholders[16].text = country
                    slide.placeholders[17].text = date.strftime("%b %Y")

                    chart_data = XyChartData()

                    cd = chart_data.add_series('Not Expensive', number_format='$#,##0')
                    for x, y in list(zip(x_values, y_values)):
                        cd.add_data_point(x, y, number_format=None)

                    cd = chart_data.add_series('Actual Price', number_format=None)
                    cd.add_data_point(Real_price, 0, number_format=None)
                    cd.add_data_point(Real_price, NE_Real_price, number_format=None)

                    cd = chart_data.add_series('Actual', number_format=None)
                    cd.add_data_point(Real_price, NE_Real_price, number_format=None)

                    graphic_frame = slide.placeholders[15].insert_chart(XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
                                                                        chart_data)
                    chart = graphic_frame.chart

                    chart.chart_style = 2  #
                    chart.has_title = False  #
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    chart.legend.font.size = Pt(12)
                    chart.legend.font.color.rgb = RGBColor(89, 89, 89)
                    if color_settings == "CMI":
                        chart.font.name = 'Rawline'
                    else:
                        chart.font.name = 'Raleway'

                    category_axis = chart.category_axis
                    category_axis.has_major_gridlines = True  #
                    category_axis.tick_labels.font.size = Pt(9)
                    category_axis.tick_labels.font.color.rgb = RGBColor(89, 89, 89)
                    category_axis.maximum_scale = maximum  # Maximum
                    category_axis.minimum_scale = minimum  # minimum
                    category_axis.major_gridlines.format.line.color.rgb = RGBColor(217, 217, 217)
                    category_axis.format.line.color.rgb = RGBColor(217, 217, 217)
                    category_axis.major_unit = units

                    value_axis = chart.value_axis
                    value_axis.has_major_gridlines = False
                    value_axis.tick_labels.font.size = Pt(9)
                    value_axis.maximum_scale = 1.0  # Maximum
                    value_axis.minimum_scale = 0.0  # minimum
                    tick_labels = value_axis.tick_labels  # Tick_Labels Control class for the vertical axis tag of Chart
                    tick_labels.number_format = '0%'  # Tag display style
                    tick_labels.font.color.rgb = RGBColor(89, 89, 89)
                    value_axis.major_gridlines.format.line.color.rgb = RGBColor(217, 217, 217)
                    value_axis.format.line.color.rgb = RGBColor(217, 217, 217)

                    plot = chart.plots[0]
                    series = plot.series[0]
                    series.format.line.color.rgb = RGBColor(28, 109, 125)
                    series.format.line.width = Pt(2.25)

                    plot = chart.plots[0]
                    series = plot.series[1]
                    series.format.line.color.rgb = RGBColor(28, 109, 125)
                    series.format.line.dash_style = 4

                    plot = chart.plots[0]
                    series = plot.series[2]
                    series.format.line.color.rgb = RGBColor(28, 109, 125)
                    series.marker.size = 8
                    series.marker.format.fill.solid()
                    series.marker.format.line.fill.background()
                    series.marker.format.fill.fore_color.rgb = RGBColor(28, 109, 125)
                    series.marker.style = XL_MARKER_STYLE.SQUARE

                    data_label = series.points[0].data_label
                    data_label.position = XL_DATA_LABEL_POSITION.LEFT
                    data_label.font.color.rgb = RGBColor(28, 109, 125)
                    data_label.font.size = Pt(12)
                    data_label.font.bold = True
                    data_label.Number_Format = '0%'

                # ----------------------------------------------------------------------------------------------------------------#
                #               Add Trial & Revenue slide for each product
                if check:
                    new_df_trial = df_trial.loc[df_trial['ProductNumber'] == current]
                    new_df_trial = new_df_trial.reset_index()
                    slide = prs.slides.add_slide(slide_layout3)
                    slide.placeholders[13].text = date.strftime("%b %Y") + ' - ' + new_df.loc[0, 'ProductName']  # Orange text in sub Header
                    if legend_settings == 'Spanish':
                        slide.placeholders[0].text = country.upper() + ' - Curvas Trial, Revenue & Profit'
                    else:
                        slide.placeholders[0].text = country.upper() + ' - Trial, Revenue & Profit Curves'
                    picture = slide.placeholders[14].insert_picture(data_file10[i])  # Product Image

                    chart_data = XyChartData()

                    x_values = new_df_trial['Price'].to_numpy()
                    y_values = new_df_trial['Share'].to_numpy()
                    y_values2 = new_df_trial['Revenue'].to_numpy()
                    y_values3 = new_df_trial['Profit'].to_numpy()

                    maxValueIndex = new_df_trial.idxmax()
                    max_trial = x_values[maxValueIndex[3]]
                    max_revenue = x_values[maxValueIndex[4]]
                    max_profit = x_values[maxValueIndex[5]]
                    max_revenue_value = y_values2[maxValueIndex[4]]
                    max_trial_value = y_values[maxValueIndex[3]]
                    max_profit_value = y_values3[maxValueIndex[5]]

                    # Rescale revenue & profit to 0-99
                    y_values2 = y_values2/(1.01*max_revenue_value)
                    y_values3 = y_values3 / (1.01*max_revenue_value)

                    minimum = min(x_values)
                    maximum = max(x_values)
                    units = round((maximum - minimum) / 14 + 0.5, 0)

                    slide.placeholders[16].text = '$ ' + str(round(max_revenue,2))
                    slide.placeholders[17].text = '$ ' + str(round(max_profit, 2))
                    slide.placeholders[18].text = '$ ' + str(round(max_trial, 2))

                    cd = chart_data.add_series('Trial', number_format='$#,##0')
                    for x, y in list(zip(x_values, y_values)):
                        cd.add_data_point(x, y, number_format=None)

                    cd = chart_data.add_series('Revenue', number_format='$#,##0')
                    for x, y in list(zip(x_values, y_values2)):
                        cd.add_data_point(x, y, number_format=None)

                    cd = chart_data.add_series('Profit', number_format='$#,##0')
                    for x, y in list(zip(x_values, y_values3)):
                        cd.add_data_point(x, y, number_format=None)

                    cd = chart_data.add_series('Max Trial', number_format=None)
                    cd.add_data_point(max_trial, 0, number_format=None)
                    cd.add_data_point(max_trial, max_trial_value, number_format=None)

                    cd = chart_data.add_series('Max Trial marker', number_format=None)
                    cd.add_data_point(max_trial, max_trial_value, number_format=None)

                    cd = chart_data.add_series('Max Revenue', number_format=None)
                    cd.add_data_point(max_revenue, 0, number_format=None)
                    cd.add_data_point(max_revenue, max_revenue_value/(1.01*max_revenue_value), number_format=None)

                    cd = chart_data.add_series('Max Revenue marker', number_format=None)
                    cd.add_data_point(max_revenue, max_revenue_value/(1.01*max_revenue_value), number_format=None)

                    cd = chart_data.add_series('Max Profit', number_format=None)
                    cd.add_data_point(max_profit, 0, number_format=None)
                    cd.add_data_point(max_profit, max_profit_value/(1.01*max_revenue_value), number_format=None)

                    cd = chart_data.add_series('Max Profit marker', number_format=None)
                    cd.add_data_point(max_profit, max_profit_value/(1.01*max_revenue_value), number_format=None)

                    graphic_frame = slide.placeholders[15].insert_chart(XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS, chart_data)

                    chart = graphic_frame.chart

                    chart.chart_style = 2  #
                    chart.has_title = False  #
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    chart.legend.font.size = Pt(12)
                    chart.legend.font.color.rgb = RGBColor(89, 89, 89)
                    if color_settings == "CMI":
                        chart.font.name = 'Rawline'
                    else:
                        chart.font.name = 'Raleway'

                    category_axis = chart.category_axis
                    category_axis.has_major_gridlines = True  #
                    category_axis.tick_labels.font.size = Pt(9)
                    category_axis.tick_labels.font.color.rgb = RGBColor(89, 89, 89)
                    category_axis.maximum_scale = round(0.8*maximum,0)  # Maximum
                    category_axis.minimum_scale = round(4*minimum,0) # minimum
                    category_axis.major_gridlines.format.line.color.rgb = RGBColor(217, 217, 217)
                    category_axis.format.line.color.rgb = RGBColor(217, 217, 217)
                    category_axis.major_unit = units

                    value_axis = chart.value_axis
                    value_axis.has_major_gridlines = False
                    value_axis.tick_labels.font.size = Pt(9)
                    value_axis.maximum_scale = 1.0  # Maximum
                    value_axis.minimum_scale = 0.0  # minimum
                    value_axis.major_unit = 0.1
                    tick_labels = value_axis.tick_labels  # Tick_Labels Control class for the vertical axis tag of Chart
                    tick_labels.number_format = '0%'  # Tag display style
                    tick_labels.font.color.rgb = RGBColor(89, 89, 89)
                    value_axis.major_gridlines.format.line.color.rgb = RGBColor(217, 217, 217)
                    value_axis.format.line.color.rgb = RGBColor(217, 217, 217)

                    plot = chart.plots[0]
                    series = plot.series[0] #Trial Curve
                    series.format.line.color.rgb = RGBColor(28,109,125)
                    series.format.line.width = Pt(2.25)

                    plot = chart.plots[0]
                    series = plot.series[1] #Revenue curve
                    series.format.line.color.rgb = RGBColor(0,176,80)
                    series.format.line.width = Pt(2.25)

                    plot = chart.plots[0]
                    series = plot.series[2] # Profit Curve
                    series.format.line.color.rgb = RGBColor(239,123,38)
                    series.format.line.width = Pt(2.25)

                    plot = chart.plots[0]
                    series = plot.series[3] #Max Trial
                    series.format.line.color.rgb = RGBColor(28,109,125)
                    series.format.line.dash_style = 4

                    plot = chart.plots[0]
                    series = plot.series[4] # Max Trial Marker
                    series.format.line.color.rgb = RGBColor(28,109,125)
                    series.marker.size = 10
                    series.marker.format.fill.solid()
                    series.marker.format.line.fill.background()
                    series.marker.format.fill.fore_color.rgb = RGBColor(28,109,125)
                    series.marker.style = XL_MARKER_STYLE.SQUARE

                    plot = chart.plots[0]
                    series = plot.series[5] # Max Revenue
                    series.format.line.color.rgb = RGBColor(0,176,80)
                    series.format.line.dash_style = 4

                    plot = chart.plots[0]
                    series = plot.series[6] # Max Revenue marker
                    series.format.line.color.rgb = RGBColor(0,176,80)
                    series.marker.size = 10
                    series.marker.format.fill.solid()
                    series.marker.format.line.fill.background()
                    series.marker.format.fill.fore_color.rgb = RGBColor(0,176,80)
                    series.marker.style = XL_MARKER_STYLE.SQUARE

                    plot = chart.plots[0]
                    series = plot.series[7] # Max profit
                    series.format.line.color.rgb = RGBColor(239,123,38)
                    series.format.line.dash_style = 4

                    plot = chart.plots[0]
                    series = plot.series[8] # Max profit marker
                    series.format.line.color.rgb = RGBColor(239,123,38)
                    series.marker.size = 10
                    series.marker.format.fill.solid()
                    series.marker.format.line.fill.background()
                    series.marker.format.fill.fore_color.rgb = RGBColor(239,123,38)
                    series.marker.style = XL_MARKER_STYLE.SQUARE

            # ----------------------------------------------------------------------------------------------------------------#
            #                   Save Presentation & open for Check
            output = BytesIO()
            prs.save(output)
            output.seek(0)

            download = FileDownloader(output.read()).download()

            st.success("PSS Report created successfully")
            # os.startfile("{}/PSS_output_NEW.pptx".format(file_path))

#-------------------------------------------In case of any issues with the sidebar---------------------------------#
else:
    st.subheader("About")
